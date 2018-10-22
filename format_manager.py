
'''
defendency
pdf - pdfminer.six
excel - openpyxl
word - python-docx
pptx - python-pptx


'''

import docx
import openpyxl
from pptx import Presentation
#pdf modules
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import  LAParams
from pdfminer.pdfpage import  PDFPage
from io import StringIO

# image change modules
from PIL import Image      # pip install pillow
from pytesseract import *  # pip install pytesseract


def image_change_format(path, lang='eng+kor'):  # path = 파일경로, lang = 적용언어(영어+한글)
    im = Image.open(path)  # path image open
    try:
        text = image_to_string(im, lang=lang, config='--psm 1 -c preserve_interword_spaces=1')
        # image -> string 변환(대상 이미지, 변환 언어, 변환 룰 설정(글자간 거리 등 *건드리지 말 것*))

    except:  # 변환되는 텍스트가 없을 경우 예외처리
        print('{} won\'t allow text extraction!'.format(path))
        text = ""
    return text






# 한글없음



#엑셀 문제있음?
def excel_change_format(path):

    wb = openpyxl.load_workbook(path)
    worksheet = wb.active

    excel_text = str()

    for col in worksheet.columns:
        max_length = 0
        column = col[0].column
        for cell in col:
            try:  # 빈셀 에러방지
                if len(str(cell.value)) > max_length:
                    max_length = len(cell.value)
                    # excel_text += ' ' + cell.value
            except:
                pass
            if cell.value is not None:
                excel_text += ' ' + str(cell.value)
        adjusted_width = (max_length + 2) * 1.2
        worksheet.column_dimensions[column].width = adjusted_width

    return excel_text


def word_change_format(path):
    #path = "C:/Users/kitri/Desktop/2012-01-11.docx"
    output_text = str()

    document = docx.Document(path)
    doc_p = ' '.join([
        paragraph.text for paragraph in document.paragraphs
    ])

    output_text += doc_p

    tables = document.tables
    for table in tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    output_text += ' '+paragraph.text

    return output_text





def pptx_change_format(path):
    #path = "C:/Users/kitri/Desktop/프로젝트관리_발표자료.pptx"
    prs = Presentation(path)

    ppt_text = str()
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                ppt_text += ' '+shape.text

    return ppt_text



def pdf_change_format(path):

    rsrcmgr = PDFResourceManager(); retstr = StringIO(); codec = 'utf-8'
    laparams = LAParams(); device = TextConverter(rsrcmgr, retstr, codec=codec, laparams=laparams)

    fp = open(path, 'rb')
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    maxpages = 0; caching = True; pagenos=set()

    try:
        for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages, caching=caching, check_extractable=True):
            interpreter.process_page(page)
    except:  # 추출되는 텍스트 없을 경우 예외처리
        print('{} won\'t allow text extraction!'.format(path))

    # for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages,caching=caching, check_extractable=True):
    # interpreter.process_page(page)

    text = retstr.getvalue()

    fp.close(); device.close(); retstr.close()

    return text

def read_txt(path):
    file = open(path, mode='r',errors='ignore')

    lines = file.readlines()

    output_text =str()

    for each_line in lines:
        output_text +=' '+each_line

    return output_text



# 프로세스 와 프로세스 경로를 주면 텍스트로 변환해서 던져주는 함수
def get_text_to_process(pname,path_list):
    output_text = str()

    if pname is 'notepad':
        if len(path_list) is not 1:
            for i in range(0, len(path_list)):
                read_path = path_list[i][1]
                output_text +=' ' + read_txt(read_path)
        else:
            for path in path_list:
                output_text += ' ' + read_txt(path[1])
        return output_text



    if pname is 'POWERPNT':
        if len(path_list) is not 1:
            for i in range(0, len(path_list)):
                read_path = path_list[i][1]
                output_text +=' ' + pptx_change_format(read_path)
        else:
            for path in path_list:
                output_text += ' ' + pptx_change_format(path[1])
        return output_text


    if pname is 'AcroRd32':
        if len(path_list) is not 1:
            for i in range(0, len(path_list)):
                read_path = path_list[i][1]
                output_text += ' ' + pdf_change_format(read_path)
        else:
            for path in path_list:
                output_text += ' ' + pdf_change_format(path[1])
        return output_text

    if pname is 'winword':
        if len(path_list) is not 1:
            for i in range(0, len(path_list)):
                read_path = path_list[i][1]
                output_text += ' ' + word_change_format(read_path)
        else:
            for path in path_list:
                output_text += ' ' + word_change_format(path[1])
        return output_text

    if pname is 'excel':
        if len(path_list) is not 1:
            for i in range(0, len(path_list)):
                read_path = path_list[i][1]
                output_text += ' ' + excel_change_format(read_path)
        else:
            for path in path_list:
                output_text += ' ' + excel_change_format(path[1])
        return output_text
    if pname is 'hwp':
        return '준비중'







